#!/usr/bin/env python3
"""
build_pitchbook_axone.py -- génère Pernod_Ricard_Pitchbook_Axone.pptx (12 slides, conventions IB).
Pre-requis : 12 images chart_01.png ... chart_12.png (graphiques Excel) dans le dossier EXH.
pip install python-pptx
"""
import os, struct
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- chemins ----------
BASE = os.path.dirname(os.path.abspath(__file__))
EXH  = os.path.abspath(os.path.join(BASE, "..", "001_pernod_ricard", "exhibits"))
OUT  = os.path.abspath(os.path.join(BASE, "..", "001_pernod_ricard", "Pernod_Ricard_Pitchbook_Axone.pptx"))

# ---------- charte Axone Equity Research ----------
NAVY  = RGBColor(0x00,0x33,0x66) # Bleu IB Classique
BLUE  = RGBColor(0x4F,0x81,0xBD) # Accentuation
GREY  = RGBColor(0x73,0x73,0x73)
LGREY = RGBColor(0xF2,0xF2,0xF2)
DARK  = RGBColor(0x26,0x26,0x26)
WHITE = RGBColor(0xFF,0xFF,0xFF)
ACCENT= RGBColor(0xD1,0xD5,0xDB)
FONT  = "Arial"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def png_size(path):
    if not os.path.exists(path): return (1600,1000) 
    with open(path,"rb") as f:
        head=f.read(24)
    if head[:8]!=b"\x89PNG\r\n\x1a\n": return (1600,1000)
    w,h=struct.unpack(">II",head[16:24]); return (w,h)

def add(): return prs.slides.add_slide(BLANK)

def set_run(r, text, size, bold=False, color=DARK, italic=False):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color

def banner(s, title):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame; tf.margin_left = Inches(0.45); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(tf.paragraphs[0].add_run(), title, 22, True, WHITE)
    
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), SW, Inches(0.06))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background(); ln.shadow.inherit = False

def footer(s, n, custom_source=None):
    # Position Y stricte et identique pour TOUTES les slides (évite l'effet de "saut")
    y_pos = Inches(7.15)
    
    # Gestion des sources spécifiques (Bloomberg, Mergermarket, etc.)
    txt = custom_source if custom_source else "Source: Company filings, Axone Equity Research estimates | Strictly Private & Confidential"
    
    # Textbox de la source (alignée à gauche)
    tb = s.shapes.add_textbox(Inches(0.45), y_pos, Inches(11), Inches(0.3))
    set_run(tb.text_frame.paragraphs[0].add_run(), txt, 8, False, GREY)
    
    # Textbox du numéro de page (alignée à droite)
    if n is not None:
        pn = s.shapes.add_textbox(Inches(12.5), y_pos, Inches(0.4), Inches(0.3))
        p = pn.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        set_run(p.add_run(), str(n), 8, False, GREY)

def fit(path, box_l, box_t, box_w, box_h):
    iw,ih = png_size(path); r = iw/ih; bw,bh = box_w,box_h
    if bw/bh > r:  w = Emu(int(bh*r)); h = Emu(int(bh))
    else:          w = Emu(int(bw));   h = Emu(int(bw/r))
    l = Emu(int(box_l+(box_w-w)/2)); t = Emu(int(box_t+(box_h-h)/2))
    return l, t, w, h

def add_safe_picture(s, path, l, t, w, h):
    """Ajoute le graphique Excel, ou un rectangle gris s'il est manquant."""
    if os.path.exists(path):
        s.shapes.add_picture(path, l, t, width=w, height=h)
    else:
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = LGREY
        box.line.color.rgb = GREY
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = f"[Graphique manquant : {os.path.basename(path)}]"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = GREY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def infobox(s, left, top, width, height, title, lines, fill=LGREY):
    # Style IB : Rectangle net avec bordure subtile (gris clair)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = ACCENT; box.line.width = Pt(1)
    box.shadow.inherit = False
    
    # Ligne supérieure bleu marine uniquement si on a un titre (très corporate)
    if title:
        top_line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.05))
        top_line.fill.solid(); top_line.fill.fore_color.rgb = NAVY; top_line.line.fill.background()

    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    
    if title:
        p_title = tf.paragraphs[0]
        set_run(p_title.add_run(), title, 12, True, NAVY)
        p_title.space_after = Pt(4)
        start_idx = 1
    else:
        start_idx = 0

    for idx, ln in enumerate(lines):
        p = tf.paragraphs[0] if start_idx == 0 and idx == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        if title or len(lines) > 1:
            set_run(p.add_run(), "▪  ", 10, True, NAVY) # Puce propre
        set_run(p.add_run(), ln, 11, False, DARK)

def four_tiles(s, items):
    """Crée 4 colonnes (tuiles) pour présenter des idées clés (Executive Summary style)."""
    left = 0.5
    top = 1.6
    w = 2.9
    h = 5.0
    gap = 0.23
    
    for i, (title, body) in enumerate(items):
        # Contour global
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + i*(w+gap)), Inches(top), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = LGREY
        box.line.color.rgb = NAVY; box.line.width = Pt(1)
        
        # En-tête bleu
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + i*(w+gap)), Inches(top), Inches(w), Inches(0.8))
        head.fill.solid(); head.fill.fore_color.rgb = NAVY
        head.line.fill.background()
        tf_h = head.text_frame; tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_h.margin_left = Inches(0.1); tf_h.margin_right = Inches(0.1)
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.alignment = PP_ALIGN.CENTER
        set_run(p_h.add_run(), title, 14, True, WHITE)
        
        # Corps de texte (centré verticalement)
        tf_b = box.text_frame; tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_b.margin_top = Inches(0.8); tf_b.margin_left = Inches(0.2); tf_b.margin_right = Inches(0.2)
        tf_b.word_wrap = True
        set_run(tf_b.paragraphs[0].add_run(), body, 13, False, DARK)

def chart_with_takeaways(s, image_filename, takeaway_items):
    """Graphique à gauche et texte centré verticalement à droite."""
    p = os.path.join(EXH, image_filename)
    l, t, w, h = fit(p, Inches(0.4), Inches(1.2), Inches(8.0), Inches(5.5))
    add_safe_picture(s, p, l, t, w, h)
    
    # Boîte de texte à droite, centrée verticalement
    tb = s.shapes.add_textbox(Inches(8.6), Inches(1.2), Inches(4.3), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE 
    
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Key Takeaways", 18, True, NAVY)
    p.space_after = Pt(12)
    
    for head, body in takeaway_items:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(8)
        p.level = 0
        set_run(p.add_run(), head + " : ", 12, True, NAVY)
        set_run(p.add_run(), body, 12, False, DARK)

def two_col_bullets(s, col1, col2, left=0.6, top=1.4, width=5.8, gap=15):
    """Puces sur 2 colonnes, centrées verticalement."""
    for col_idx, items in enumerate([col1, col2]):
        pos_left = left if col_idx == 0 else left + width + 0.5
        tb = s.shapes.add_textbox(Inches(pos_left), Inches(top), Inches(width), Inches(5.2))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, (head, body) in enumerate(items):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            set_run(p.add_run(), "▪  " + head + " : ", 14, True, NAVY)
            set_run(p.add_run(), body, 14, False, DARK)

# ============================================================
# SLIDE 1 — COVER
# ============================================================
s=add()
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(2.3),SW,Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb=NAVY; band.line.fill.background(); band.shadow.inherit=False
tf=band.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; set_run(p.add_run(), "Pernod Ricard SA  (RI FP)", 40, True, WHITE)
p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; set_run(p.add_run(), "Initiating Coverage  ·  Spirits & Wine", 18, False, ACCENT)
p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; p.space_before=Pt(10)
set_run(p.add_run(), "BUY   ·   Target €108   ·   Upside +68%   ·   Current €64.14", 22, True, WHITE)

# Footer de la Cover aligné exactement au même niveau que les autres slides (7.15)
tb=s.shapes.add_textbox(0, Inches(7.15), SW, Inches(0.3))
p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
# Mention discrète Sabashvili Rezi
set_run(p.add_run(), "Axone Equity Research   |   Prepared by Sabashvili Rezi   |   Juin 2026", 11, False, GREY)

# ============================================================
# SLIDE 2 — INVESTMENT THESIS
# ============================================================
s=add(); banner(s, "Investment thesis: a quality franchise at a discount")
four_tiles(s, [
 ("Deep valuation discount", "Pernod trades at 8.1x EV/EBITDA versus a 12.4x peer median. This implies a ~24% discount despite holding the #2 global position and an ultra-premium portfolio (Jameson, Absolut, Martell)."),
 ("FY26 transition year", "Q3 organic net sales returned to stability (+0.1%) after a soft H1. Management guidance of +3% to +6% organic growth per annum over FY27-29 offers clear visibility."),
 ("Clear deleveraging", "We project net debt/EBITDA falling from ~3.4x to ~2.0x by FY30. This is perfectly in line with the <3.0x target by FY29, supported by an ~80% cash conversion rate."),
 ("M&A floors valuation", "Precedent transactions in the spirits sector have priced at a median 16.9x EV/EBITDA. A control-premium scenario yields an implied value well above our €108 target."),
])
footer(s, 2)

# ============================================================
# SLIDE 3 — COMPANY OVERVIEW
# ============================================================
s=add(); banner(s, "A diversified, premium global spirits leader")
chart_with_takeaways(s, "chart_09.png", [
    ("Global Footprint", "Revenues are highly diversified across the Americas, Europe, and Asia/RoW, minimizing reliance on a single geographic market."),
    ("Premiumization Strategy", "High exposure to premium and super-premium segments protects gross margins during inflationary cycles."),
    ("Strategic Focus", "Continued rationalization of the portfolio (e.g., disposal of local wine brands) to focus exclusively on high-margin core spirits.")
])
footer(s, 3)

# ============================================================
# SLIDE 4 — FINANCIAL SUMMARY
# ============================================================
s=add(); banner(s, "Returning to top-line growth with steady margin expansion")
# Graphiques jumeaux
for fn,bl in [("chart_07.png",Inches(0.4)),("chart_08.png",Inches(6.8))]:
    p_img=os.path.join(EXH,fn); l,t,w,h=fit(p_img,bl,Inches(1.15),Inches(6.1),Inches(4.8))
    add_safe_picture(s, p_img, l, t, w, h)
# Remontée de l'infobox pour ne pas cacher le footer
infobox(s, 0.4, 6.2, 12.5, 0.65, "", [
    "Takeaway: We model organic revenue growth stabilizing at ~4-5% from FY28, with gross margins expanding slightly to 60.0% by FY30 driven by positive price-mix and easing COGS."
], fill=WHITE)
footer(s, 4)

# ============================================================
# SLIDE 5 — VALUATION SUMMARY
# ============================================================
s=add(); banner(s, "Valuation points to material upside across all methodologies")
chart_with_takeaways(s, "chart_01.png", [
    ("DCF Valuation", "Points to a fundamental range of €88 - €99 per share, based on conservative terminal growth (1.5%) and exit multiples (11.0x)."),
    ("Trading Comparables", "Suggests a fair value of ~€105. Pernod's current multiple is heavily penalized compared to direct peers like Diageo and Campari."),
    ("Precedent Transactions", "Implies a control value of ~€161, highlighting the immense strategic scarcity value of Pernod's portfolio."),
    ("Conclusion", "Our €108 blended target price represents an asymmetric risk/reward profile with +68% upside from current spot levels.")
])
footer(s, 5)

# ============================================================
# SLIDE 6 — DCF ANALYSIS
# ============================================================
s=add(); banner(s, "DCF supports a €88-99 intrinsic value")
for fn,bl in [("chart_05.png",Inches(0.4)),("chart_06.png",Inches(6.8))]:
    # Réduction légère de la hauteur des graphiques (4.1 au lieu de 4.2) pour laisser respirer le texte
    p_img=os.path.join(EXH,fn); l,t,w,h=fit(p_img,bl,Inches(1.25),Inches(6.1),Inches(4.1))
    add_safe_picture(s, p_img, l, t, w, h)

# Encadré IB propre, plus large, mieux proportionné, avec des arguments scindés
infobox(s, 0.4, 5.5, 12.5, 1.4, "Key Modeling Assumptions", [
    "Discount Rate & TV: WACC of 7.3% (Beta 1.25, ERP 5.5%, RFR 3.7%). Terminal growth of 1.5% and implied exit multiple of 11.0x.",
    "Cash Flow Mechanics: Unlevered FCF = NOPAT + D&A - Capex - Change in NWC. Effective tax rate normalized at 25%.",
    "Valuation Triangulation: Enterprise Value derived from an equal blend of Gordon Growth (€34.4bn) and Exit Multiple (€37.1bn) methods."
])
footer(s, 6)

# ============================================================
# SLIDE 7 — DCF SENSITIVITY
# ============================================================
s=add(); banner(s, "Valuation is robust across WACC and growth scenarios")
chart_with_takeaways(s, "chart_12.png", [
    ("Base Case", "At a 7.3% WACC and 1.5% perpetuity growth, the Gordon Growth method yields €88/share."),
    ("Downside Protection", "Even in a stressed scenario (WACC at 8.0%, growth at 1.0%), the implied equity value remains above the current market price of €64.14."),
    ("Upside Potential", "A reversion to historical WACC levels (~6.5%) combined with standard 2.0% terminal growth would imply a share price above €120.")
])
footer(s, 7)

# ============================================================
# SLIDE 8 — TRADING COMPS
# ============================================================
s=add(); banner(s, "Pernod is mispriced versus global spirits peers")
for fn,bl in [("chart_02.png",Inches(0.4)),("chart_03.png",Inches(6.8))]:
    p_img=os.path.join(EXH,fn); l,t,w,h=fit(p_img,bl,Inches(1.15),Inches(6.1),Inches(4.8))
    add_safe_picture(s, p_img, l, t, w, h)
infobox(s, 0.4, 6.2, 12.5, 0.65, "", [
    "Takeaway: At 8.1x EV/EBITDA, PR trades at a ~35% discount to Diageo (10.8x) and Campari (12.8x). Reverting to the peer median (12.4x) yields our €105 implied price."
], fill=WHITE)
footer(s, 8, "Source: Bloomberg, Company filings, Axone Equity Research estimates | Strictly Private & Confidential")

# ============================================================
# SLIDE 9 — PRECEDENT TRANSACTIONS
# ============================================================
s=add(); banner(s, "M&A precedents imply a substantial control premium")
chart_with_takeaways(s, "chart_04.png", [
    ("Historical Multiples", "Over the last 15 years, major spirits platforms (Beam, Allied Domecq, Grand Marnier) have been acquired at a median EV/EBITDA of 16.9x."),
    ("Strategic Scarcity", "Assets of Pernod Ricard's scale are essentially impossible to replicate. Any theoretical buyout would require a massive premium."),
    ("Implied Value", "Applying the 16.9x median to our FY26E EBITDA (€3.1bn) yields an implied share price of ~€161, acting as a definitive valuation floor.")
])
footer(s, 9, "Source: Mergermarket, Company filings, Axone Equity Research estimates | Strictly Private & Confidential")

# ============================================================
# SLIDE 10 — VALUATION TRIANGULATION
# ============================================================
s=add(); banner(s, "Triangulating to a €108 blended target price")
chart_with_takeaways(s, "chart_11.png", [
    ("Methodology Weighting", "We apply a 50% weight to DCF (average of Gordon and Exit methods), 35% to Trading Comps, and 15% to Precedent Transactions."),
    ("Why this mix?", "Acknowledges the fundamental cash-generative nature of the business (DCF) while factoring in the reality of public market multiples (Comps) and latent strategic value (M&A)."),
    ("Final Target", "The blended output is €107.91, which we round to €108, representing a 68% upside.")
])
footer(s, 10)

# ============================================================
# SLIDE 11 — CAPITAL STRUCTURE
# ============================================================
s=add(); banner(s, "Strong cash generation drives rapid deleveraging by FY29")
chart_with_takeaways(s, "chart_10.png", [
    ("Current Position", "Net debt stands at €11.1bn (FY25), representing an elevated ~3.4x Net Debt/EBITDA ratio."),
    ("Cash Conversion", "Robust Unlevered FCF generation allows for rapid debt paydown without compromising the dividend policy (DPS €4.70)."),
    ("Target Reached", "We forecast leverage falling comfortably below management's 3.0x target by FY29, and reaching 2.0x by FY30, offering future balance sheet optionality.")
])
footer(s, 11)

# ============================================================
# SLIDE 12 — RISKS & DISCLAIMER
# ============================================================
s=add(); banner(s, "Key risks to our thesis")
col1 = [
 ("China weakness", "Organic net sales -24% YTD; prolonged Martell/cognac softness would severely pressure margin mix."),
 ("US normalization", "Organic net sales -14% YTD. The pace of trade inventory adjustments taking longer than expected could push the return to growth into late FY27."),
 ("Adverse FX", "Management guides to a significantly negative FX impact (USD, INR, TRY) presenting a material translation drag on EUR earnings.")
]
col2 = [
 ("Global Travel Retail", "Growth could stall due to geopolitical tensions and Middle East travel disruptions, despite a strong Q3 (+39%)."),
 ("Execution risk", "The €1bn FY26-29 efficiency program is central to margin defense. Any under-delivery directly hits our DCF value."),
 ("Tariff escalations", "Potential trade disputes (e.g., EU vs China cognac probes) could result in punitive tariffs that compress margins.")
]
# Puces réparties sur deux colonnes
two_col_bullets(s, col1, col2)

db=s.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(12), Inches(0.5))
db.text_frame.word_wrap=True
set_run(db.text_frame.paragraphs[0].add_run(),
        "Educational purposes only; not investment advice, an offer, or a solicitation. "
        "Prepared for the Axone platform based on public filings and proprietary estimates. Figures as of June 2026.",
        8, False, GREY, italic=True)
footer(s, 12)

prs.save(OUT)
print("Génération réussie ->", OUT)